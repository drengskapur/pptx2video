#!/usr/bin/env python3
"""
.. module:: scripts.pptx2video
   :synopsis: Convert PowerPoint presentations to video with AI-generated voiceovers.

This script converts a PowerPoint presentation (.pptx) into a video (.mp4) by:
- Extracting slides as images (via PDF intermediate).
- Generating voiceover audio using OpenAI TTS API based on slide notes.
- Combining images and audio into video segments per slide.
- Concatenating the slide videos into a final video.

The process is idempotent and utilizes both asynchronous and parallel processing
to improve performance.

Usage:
    python pptx2video.py <presentation.pptx> [--output <output.mp4>] [--speed <speed>]

Requirements:
    - Python 3.7+
    - Install required packages:
        pip install python-pptx Pillow openai pdf2image tenacity
    - LibreOffice must be installed and accessible in the system PATH.
    - FFmpeg must be installed and accessible in the system PATH.
    - Set the 'OPENAI_API_KEY' environment variable with your OpenAI API key.

Note:
    Ensure you have sufficient permissions and API quota for OpenAI TTS API.
"""

import argparse
import asyncio
import hashlib
import json
import logging
import os
import subprocess
from typing import List, Optional
import multiprocessing
from concurrent.futures import ProcessPoolExecutor, ThreadPoolExecutor
import functools

import openai
from pdf2image import convert_from_path
from pptx import Presentation
from PIL import Image
from tenacity import retry, stop_after_attempt, wait_exponential, retry_if_exception_type

# Configure logging
logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")
logger = logging.getLogger(__name__)

# Maximum number of concurrent API calls to OpenAI
MAX_CONCURRENT_CALLS = 5


@retry(
    stop=stop_after_attempt(3),
    wait=wait_exponential(multiplier=1, min=4, max=10),
    retry=retry_if_exception_type(openai.OpenAIError),
)
async def tts_async(
    input_text: str, model: str, voice: str, api_key: str, speed: float
) -> bytes:
    """
    Asynchronously convert input text to speech using OpenAI's Text-to-Speech API.

    :param input_text: The text to be converted to speech.
    :type input_text: str
    :param model: The OpenAI TTS model to use.
    :type model: str
    :param voice: The voice profile to use.
    :type voice: str
    :param api_key: OpenAI API key.
    :type api_key: str
    :param speed: The speed of the synthesized speech (0.25 to 4.0).
    :type speed: float
    :returns: Audio content as bytes.
    :rtype: bytes
    :raises ValueError: If API key or input text is empty.
    :raises RuntimeError: If the API call fails after retries.
    """
    if not api_key.strip() or not input_text.strip():
        raise ValueError("API key and input text are required.")

    logger.debug(f"Sending TTS request for text: {input_text[:50]}...")

    with ThreadPoolExecutor() as executor:
        loop = asyncio.get_event_loop()
        try:
            speech_file = await loop.run_in_executor(
                executor,
                functools.partial(
                    openai.audio.speech.create,
                    model=model.lower(),
                    voice=voice.lower(),
                    input=input_text,
                    response_format="mp3",
                    speed=speed,
                    api_key=api_key,
                ),
            )
        except openai.OpenAIError as e:
            logger.error(f"OpenAI API error: {e}")
            raise  # Re-raise for retry mechanism
        except Exception as e:
            logger.error(f"Unexpected error during TTS generation: {e}")
            raise

    logger.debug(f"Received audio content, size: {len(speech_file.content)} bytes")
    return speech_file.content


@retry(stop=stop_after_attempt(3), wait=wait_exponential(multiplier=1, min=4, max=10))
def convert_slide_to_image(pdf_filename: str, output_dir: str, i: int) -> str:
    """
    Convert a single slide from the PDF to an image.

    :param pdf_filename: Path to the PDF file.
    :type pdf_filename: str
    :param output_dir: Directory to store output files.
    :type output_dir: str
    :param i: Slide index.
    :type i: int
    :returns: Path to the generated image file.
    :rtype: str
    """
    logger.debug(f"Converting slide {i+1} to image")
    image_path = os.path.join(output_dir, f"slide_{i}.png")
    try:
        images = convert_from_path(pdf_filename, first_page=i + 1, last_page=i + 1, dpi=300)
        if images:
            images[0].save(image_path, "PNG")
            logger.debug(f"Saved slide {i+1} image to {image_path}")
        else:
            logger.warning(f"Failed to convert slide {i+1} to image")
    except Exception as e:
        logger.error(f"Error converting slide {i+1} to image: {e}")
        raise  # Re-raise for retry mechanism
    return image_path


async def generate_audio_for_slide(
    text: str, output_dir: str, i: int, api_key: str, speed: float
) -> Optional[str]:
    """
    Generate audio for a single slide using the TTS function.

    :param text: Text to convert to speech.
    :type text: str
    :param output_dir: Directory to store output files.
    :type output_dir: str
    :param i: Slide index.
    :type i: int
    :param api_key: OpenAI API key.
    :type api_key: str
    :param speed: The speed of the synthesized speech (0.25 to 4.0).
    :type speed: float
    :returns: Path to the generated audio file, or None if generation fails.
    :rtype: Optional[str]
    """
    slide_audio_filename = os.path.join(output_dir, f"voice_{i}.mp3")
    try:
        logger.debug(f"Generating audio for slide {i+1} with text: {text[:50]}...")
        audio_content = await tts_async(text, "tts-1-hd", "echo", api_key, speed)
        with open(slide_audio_filename, "wb") as f:
            f.write(audio_content)
        logger.debug(f"Generated audio for slide {i+1}: {slide_audio_filename}")
        return slide_audio_filename
    except Exception as e:
        logger.error(f"Error generating audio for slide {i+1}: {e}")
        return None


@retry(stop=stop_after_attempt(3), wait=wait_exponential(multiplier=1, min=4, max=10))
def create_video_for_slide(
    image_file: str, audio_file: Optional[str], output_dir: str, i: int
) -> str:
    """
    Create a video for a single slide by combining image and audio.

    :param image_file: Path to the slide image file.
    :type image_file: str
    :param audio_file: Path to the audio file, or None if no audio.
    :type audio_file: Optional[str]
    :param output_dir: Directory to store output files.
    :type output_dir: str
    :param i: Slide index.
    :type i: int
    :returns: Path to the generated video file.
    :rtype: str
    """
    logger.debug(f"Creating video for slide {i+1}")
    slide_video_filename = os.path.join(output_dir, f"video_{i}.mp4")
    try:
        with Image.open(image_file) as img:
            width, height = img.size
        adjusted_width = width if width % 2 == 0 else width - 1

        if audio_file:
            ffmpeg_cmd = [
                "ffmpeg",
                "-y",
                "-loop",
                "1",
                "-i",
                image_file,
                "-i",
                audio_file,
                "-c:v",
                "libx264",
                "-tune",
                "stillimage",
                "-c:a",
                "aac",
                "-b:a",
                "192k",
                "-pix_fmt",
                "yuv420p",
                "-vf",
                f"scale={adjusted_width}:-2",
                "-shortest",
                slide_video_filename,
            ]
        else:
            logger.warning(f"No audio file for slide {i+1}, creating silent video")
            ffmpeg_cmd = [
                "ffmpeg",
                "-y",
                "-loop",
                "1",
                "-i",
                image_file,
                "-c:v",
                "libx264",
                "-t",
                "5",
                "-pix_fmt",
                "yuv420p",
                "-vf",
                f"scale={adjusted_width}:-2",
                slide_video_filename,
            ]

        subprocess.run(ffmpeg_cmd, check=True, capture_output=True)
        logger.debug(f"Created video for slide {i+1}: {slide_video_filename}")
    except Exception as e:
        logger.error(f"Error creating video for slide {i+1}: {e}")
        raise  # Re-raise for retry mechanism
    return slide_video_filename


class PPTXtoVideo:
    """
    Converts a PowerPoint presentation to a video.

    :param pptx_filename: The path to the PowerPoint (.pptx) file.
    :type pptx_filename: str
    :param output_file: The path to the output video file (optional).
    :type output_file: str, optional
    """

    def __init__(self, pptx_filename: str, output_file: Optional[str] = None):
        """Initialize the PPTXtoVideo instance."""
        self.pptx_filename = pptx_filename
        self.pptx_hash = self._compute_file_hash(pptx_filename)
        self.output_dir = os.path.join(os.getcwd(), "video-assets-current")
        os.makedirs(self.output_dir, exist_ok=True)
        self.pdf_filename = os.path.join(
            self.output_dir,
            f"{os.path.splitext(os.path.basename(pptx_filename))[0]}.pdf",
        )

        # Set output file name based on input or default
        if output_file:
            self.output_file = output_file
        else:
            self.output_file = os.path.join(
                os.path.dirname(pptx_filename),
                f"{os.path.splitext(os.path.basename(pptx_filename))[0]}.mp4",
            )

        try:
            self.presentation = Presentation(pptx_filename)
        except Exception as e:
            logger.error(f"Error opening PowerPoint presentation: {e}")
            raise
        self.slides = self.presentation.slides
        self.api_key = os.environ.get("OPENAI_API_KEY")
        if not self.api_key:
            raise ValueError("OPENAI_API_KEY not found in environment variables.")
        self.voiceover_texts = self._extract_slide_texts()
        self.video_files = []
        self.state_file = os.path.join(
            self.output_dir, f"conversion_state_{self.pptx_hash[:8]}.json"
        )
        self.state = self._load_state()

    def _compute_file_hash(self, filename: str) -> str:
        """
        Compute the SHA256 hash of a file.

        :param filename: Path to the file.
        :type filename: str
        :returns: SHA256 hash of the file contents.
        :rtype: str
        """
        hasher = hashlib.sha256()
        with open(filename, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hasher.update(chunk)
        return hasher.hexdigest()

    def _extract_slide_texts(self) -> List[str]:
        """
        Extract text from slide notes.

        :returns: List of texts extracted from slide notes.
        :rtype: List[str]
        """
        texts = []
        for i, slide in enumerate(self.slides):
            text = (
                slide.notes_slide.notes_text_frame.text.strip()
                if slide.has_notes_slide
                else ""
            )
            texts.append(text)
            logger.debug(f"Slide {i+1} text: {text[:50]}...")
        return texts

    def _load_state(self) -> dict:
        """
        Load the conversion state from a file.

        :returns: Dictionary containing the conversion state.
        :rtype: dict
        """
        if os.path.exists(self.state_file):
            try:
                with open(self.state_file, "r") as f:
                    return json.load(f)
            except Exception as e:
                logger.warning(f"Error loading state file: {e}. Starting fresh.")
        return {
            "pdf_created": False,
            "images_created": [],
            "audio_created": [],
            "videos_created": [],
        }

    def _save_state(self):
        """Save the current conversion state to a file."""
        try:
            with open(self.state_file, "w") as f:
                json.dump(self.state, f)
        except Exception as e:
            logger.warning(f"Error saving state file: {e}")

    def _convert_to_pdf(self):
        """Convert the .pptx file to a .pdf file using LibreOffice."""
        if self.state["pdf_created"] and os.path.exists(self.pdf_filename):
            logger.info(f"PDF already exists: {self.pdf_filename}")
            return

        logger.info("Converting PPTX to PDF")
        os.makedirs(os.path.dirname(self.pdf_filename), exist_ok=True)

        cmd = [
            "libreoffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            os.path.dirname(self.pdf_filename),
            self.pptx_filename,
        ]
        result = subprocess.run(cmd, check=True, capture_output=True, text=True)
        logger.debug(f"LibreOffice conversion output: {result.stdout}")
        logger.debug(f"LibreOffice conversion errors: {result.stderr}")

        if not os.path.exists(self.pdf_filename):
            raise RuntimeError(f"Failed to create PDF file: {self.pdf_filename}")

        self.state["pdf_created"] = True
        self._save_state()
        logger.info(f"PDF created successfully: {self.pdf_filename}")

    async def create_videos(self, speed: float = 1.0):
        """
        Create individual video files for each slide.

        Combines slide images and TTS audio.

        :param speed: The speed of the synthesized speech (0.25 to 4.0).
        :type speed: float
        """
        self._convert_to_pdf()

        logger.info("Converting PDF slides to images")
        image_files = []
        with multiprocessing.Pool() as pool:
            convert_func = functools.partial(
                convert_slide_to_image, self.pdf_filename, self.output_dir
            )
            for i, image_path in enumerate(pool.imap(convert_func, range(len(self.slides)))):
                if i not in self.state["images_created"]:
                    image_files.append(image_path)
                    self.state["images_created"].append(i)
                    self._save_state()
        logger.info(f"Created {len(image_files)} new slide images")

        logger.info("Generating audio for slides")
        audio_files = []

        async def generate_all_audio():
            semaphore = asyncio.Semaphore(MAX_CONCURRENT_CALLS)

            async def bounded_generate(text, i):
                if i in self.state["audio_created"]:
                    return os.path.join(self.output_dir, f"voice_{i}.mp3")
                if not text.strip():
                    logger.warning(
                        f"Skipping audio generation for slide {i+1} due to empty text"
                    )
                    return None
                async with semaphore:
                    audio_file = await generate_audio_for_slide(
                        text, self.output_dir, i, self.api_key, speed
                    )
                    if audio_file:
                        self.state["audio_created"].append(i)
                        self._save_state()
                    return audio_file

            return await asyncio.gather(
                *[bounded_generate(text, i) for i, text in enumerate(self.voiceover_texts)]
            )

        audio_files = await generate_all_audio()
        logger.info(f"Generated {len([a for a in audio_files if a])} audio files")

        logger.info("Creating videos for each slide")
        max_workers = min(multiprocessing.cpu_count(), len(self.slides))
        with ProcessPoolExecutor(max_workers=max_workers) as executor:
            futures = []
            for i, (image_file, audio_file) in enumerate(zip(image_files, audio_files)):
                if i not in self.state["videos_created"]:
                    futures.append(
                        executor.submit(
                            create_video_for_slide,
                            image_file,
                            audio_file,
                            self.output_dir,
                            i,
                        )
                    )

            for i, future in enumerate(futures):
                video_file = future.result()
                if video_file:
                    self.video_files.append(video_file)
                    self.state["videos_created"].append(i)
                    self._save_state()

        logger.info(f"Created {len(self.video_files)} individual slide videos")

    def combine_videos(self):
        """Concatenate individual slide videos into a final video."""
        if os.path.exists(self.output_file):
            logger.info(f"Final video already exists: {self.output_file}")
            return

        logger.info("Combining individual slide videos into final video")
        list_file = os.path.join(self.output_dir, "videos.txt")
        with open(list_file, "w") as f:
            for video_file in self.video_files:
                f.write(f"file '{video_file}'\n")

        ffmpeg_cmd = [
            "ffmpeg",
            "-y",
            "-f",
            "concat",
            "-safe",
            "0",
            "-i",
            list_file,
            "-c",
            "copy",
            self.output_file,
        ]
        subprocess.run(ffmpeg_cmd, check=True, capture_output=True)
        logger.info(f"Final video created: {self.output_file}")

    async def convert(self, speed: float = 1.0):
        """
        Convert the PowerPoint presentation to a video.

        :param speed: The speed of the synthesized speech (0.25 to 4.0).
        :type speed: float
        """
        logger.info(f"Starting conversion of {self.pptx_filename}")
        await self.create_videos(speed)
        self.combine_videos()
        logger.info(f"Video created successfully: {self.output_file}")


async def main():
    """Main function to parse arguments and execute conversion."""
    parser = argparse.ArgumentParser(
        description="Convert a PowerPoint presentation to a video using OpenAI TTS and FFmpeg."
    )
    parser.add_argument(
        "pptx", type=str, help="The path to the PowerPoint (.pptx) file to convert."
    )
    parser.add_argument(
        "--output",
        "-o",
        type=str,
        help="The path to the output video file (optional). Defaults to the input filename with .mp4 extension.",
    )
    parser.add_argument(
        "--speed",
        type=float,
        default=1.0,
        help="The speed of the synthesized speech (0.25 to 4.0). Defaults to 1.0.",
    )
    args = parser.parse_args()

    try:
        logger.info(f"Starting conversion process for: {args.pptx}")
        converter = PPTXtoVideo(args.pptx, args.output)
        await converter.convert(args.speed)
        logger.info("Conversion process completed successfully")
    except Exception as e:
        logger.exception(f"An unexpected error occurred: {e}")
        exit(1)


if __name__ == "__main__":
    asyncio.run(main())
